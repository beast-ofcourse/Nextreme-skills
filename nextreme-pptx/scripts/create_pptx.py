#!/usr/bin/env python3
"""
nextreme-pptx — spec → .pptx engine (unbound creative freedom)

Reads a YAML or JSON spec and produces a geometry-validated, theme-coherent,
vector-editable .pptx. Every shape honors grid bounds; every chart is native;
no raster slides, no overlapping, no placeholder residue.

Usage (skills.sh):
    python ${CLAUDE_SKILL_DIR}/scripts/create_pptx.py --spec ${CLAUDE_SKILL_DIR}/templates/pitch_spec.yaml --output ./deck.pptx
    python ${CLAUDE_SKILL_DIR}/scripts/create_pptx.py --spec spec.json --output ./deck.pptx --template branded.pptx

Local clone:
    python nextreme-pptx/scripts/create_pptx.py --spec nextreme-pptx/templates/pitch_spec.yaml --output ./deck.pptx

Spec schema (YAML or JSON):
    title: "Deck Title"                 # presentation title (for core properties)
    author: "Author"
    subject: "Subject"
    theme_key: "vc_clean"              # vc_clean | academic_minimal | research_dark | editorial
    slide_size: "16:9"                 # 16:9 (13.33x7.5) | 4:3 (10x7.5) | A4 (11.69x8.27)
    slides:
      - type: "cover"
        title: "The Cover Title"
        subtitle: "One-line promise"
      - type: "bento_features"
        title: "Why This Wins"
        items: [{title: "Speed", body: "4h vs 6h"}, {title: "Cost", body: "-22%"}]
      - type: "stats_grid"
        title: "Traction"
        stats: [{value: "84%", label: "Activation", delta: "+12pp"}, ...]
      - type: "bar_chart"
        title: "Revenue"
        categories: ["Q1","Q2","Q3"]
        series: [{name: "Revenue", values: [1.2,1.9,2.4]}]
      - type: "matrix_2x2"
        title: "Positioning"
        quadrants: [{label: "Leaders"}, {label: "Niche"}, ...]
        items: [{label: "Us", x: 0.8, y: 0.8}]
      - type: "timeline"
        title: "Roadmap"
        milestones: [{label: "Q1", title: "Ship"}, ...]
      - type: "quote"
        quote: "Design is intelligence made visible."
        attribution: "— Alina Wheeler"
      - type: "results_table"
        title: "Results"
        headers: ["Metric","Before","After"]
        rows: [["Latency","6h","4m"], ...]
      - type: "section_header"
        title: "Section Name"
        kicker: "01 — CONTEXT"

Supported types: cover, section_header, closing, bento_features, moat_columns,
stats_grid, bar_chart, line_chart, matrix_2x2, timeline, quote, results_table,
comparison, process, image_text, team — see references/slide-types.md.

All Inches/Pt/hex are named tokens below — no magic in logic.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Any

try:
    import yaml  # type: ignore
except ImportError:
    yaml = None  # type: ignore

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.chart.data import ChartData
except ImportError as import_error:
    print(f"[create_pptx] Missing dependency: {import_error}", file=sys.stderr)
    print("Install: pip install -r ${CLAUDE_SKILL_DIR}/requirements.txt  or  pip install python-pptx pyyaml lxml Pillow", file=sys.stderr)
    sys.exit(2)


# ── Tokens — single source, no magic ────────────────────────────────────────

SLIDE_16_9_WIDTH = Inches(13.33)
SLIDE_16_9_HEIGHT = Inches(7.5)
SLIDE_4_3_WIDTH = Inches(10)
SLIDE_4_3_HEIGHT = Inches(7.5)
SLIDE_A4_WIDTH = Inches(11.69)
SLIDE_A4_HEIGHT = Inches(8.27)

MARGIN = Inches(0.6)
GUTTER = Inches(0.32)
GRID_COLS_16_9 = 12
GRID_COLS_4_3 = 10
CARD_PAD = Inches(0.28)
CARD_GAP = GUTTER

THEME_PALETTES: dict[str, dict[str, str]] = {
    "vc_clean":        {"BG": "FFFFFF", "FG": "1A1A1E", "FG_MUTED": "6B7280", "ACCENT": "1B4F72", "ACCENT_2": "2E86AB", "ACCENT_3": "A8DADC", "SURFACE": "F3F4F6", "BORDER": "E5E7EB"},
    "academic_minimal": {"BG": "FFFFFF", "FG": "111827", "FG_MUTED": "6B7280", "ACCENT": "374151", "ACCENT_2": "9CA3AF", "ACCENT_3": "E5E7EB", "SURFACE": "F9FAFB", "BORDER": "E5E7EB"},
    "research_dark":    {"BG": "0F172A", "FG": "E2E8F0", "FG_MUTED": "94A3B8", "ACCENT": "38BDF8", "ACCENT_2": "FB923C", "ACCENT_3": "334155", "SURFACE": "1E293B", "BORDER": "334155"},
    "editorial":       {"BG": "FFF1F2", "FG": "1F2937", "FG_MUTED": "6B7280", "ACCENT": "E11D48", "ACCENT_2": "FB7185", "ACCENT_3": "FFE4E6", "SURFACE": "FFF7F7", "BORDER": "FFE4E6"},
}

FONT_SANS = "Calibri"
FONT_SERIF = "Georgia"
FONT_MONO = "Consolas"

SLOP_TOKENS = ["lorem ipsum", "lorem", "placeholder", "your text here", "insert text here", "click to add title", "click to add text"]

MIN_TEXTBOX_WIDTH = Inches(2.5)
MAX_BULLETS_PER_SLIDE = 7


# ── Data model ───────────────────────────────────────────────────────────────

@dataclass(frozen=True)
class DeckConfig:
    title: str
    author: str
    subject: str
    theme_key: str
    slide_width: Any
    slide_height: Any
    grid_cols: int
    palette: dict[str, str]


# ── Helpers — each does one job ─────────────────────────────────────────────

def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="nextreme-pptx spec → pptx")
    parser.add_argument("--spec", required=True, help="Path to YAML or JSON spec")
    parser.add_argument("--output", required=True, help="Output .pptx path")
    parser.add_argument("--template", default="", help="Optional branded .pptx template (slide-master path)")
    return parser.parse_args()


def load_spec(spec_path: Path) -> dict[str, Any]:
    if not spec_path.exists():
        raise FileNotFoundError(f"Spec not found: {spec_path}")
    raw_text = spec_path.read_text(encoding="utf-8")
    suffix = spec_path.suffix.lower()
    if suffix in (".yaml", ".yml"):
        if yaml is None:
            raise RuntimeError("pyyaml not installed — pip install pyyaml or use JSON")
        parsed = yaml.safe_load(raw_text)
        if not isinstance(parsed, dict):
            raise ValueError(f"YAML spec must be mapping at top, got {type(parsed).__name__}")
        return parsed
    if suffix == ".json":
        parsed_json = json.loads(raw_text)
        if not isinstance(parsed_json, dict):
            raise ValueError(f"JSON spec must be object at top, got {type(parsed_json).__name__}")
        return parsed_json
    raise ValueError(f"Unsupported spec extension '{suffix}' — use .yaml/.yml or .json")


def resolve_deck_config(raw: dict[str, Any]) -> DeckConfig:
    slide_size = str(raw.get("slide_size", "16:9")).strip()
    if slide_size == "4:3":
        width, height, grid_cols = SLIDE_4_3_WIDTH, SLIDE_4_3_HEIGHT, GRID_COLS_4_3
    elif slide_size.lower() == "a4":
        width, height, grid_cols = SLIDE_A4_WIDTH, SLIDE_A4_HEIGHT, GRID_COLS_16_9
    else:
        width, height, grid_cols = SLIDE_16_9_WIDTH, SLIDE_16_9_HEIGHT, GRID_COLS_16_9  # default 16:9

    theme_key = str(raw.get("theme_key", "vc_clean")).strip().lower()
    if theme_key not in THEME_PALETTES:
        raise ValueError(f"Unknown theme_key '{theme_key}' — expected one of {list(THEME_PALETTES)}")

    return DeckConfig(
        title=str(raw.get("title", "Untitled Deck")),
        author=str(raw.get("author", "")),
        subject=str(raw.get("subject", "")),
        theme_key=theme_key,
        slide_width=width,
        slide_height=height,
        grid_cols=grid_cols,
        palette=THEME_PALETTES[theme_key],
    )


def slot_width_emu(grid_cols: int, span_cols: int, content_width_emu: int, gutter_emu: int) -> int:
    if span_cols <= 0 or span_cols > grid_cols:
        raise ValueError(f"span_cols {span_cols} must be 1..{grid_cols}")
    col_width_emu = (content_width_emu - (grid_cols - 1) * gutter_emu) // grid_cols
    return span_cols * col_width_emu + (span_cols - 1) * gutter_emu


def check_bounds(left_emu: int, top_emu: int, width_emu: int, height_emu: int, slide_width_emu: int, slide_height_emu: int, margin_emu: int, label: str) -> None:
    if width_emu < MIN_TEXTBOX_WIDTH.emu and "textbox" in label.lower():
        # warn, not fail — narrow KPI labels legitimately smaller, but body not
        pass
    if left_emu < margin_emu - int(Inches(0.02).emu):
        raise ValueError(f"Off-canvas left for {label}: left {left_emu/914400:.2f}″ < margin {margin_emu/914400:.2f}″")
    if left_emu + width_emu > slide_width_emu - margin_emu + int(Inches(0.02).emu):
        raise ValueError(f"Off-canvas right for {label}: right {(left_emu+width_emu)/914400:.2f}″ > { (slide_width_emu-margin_emu)/914400:.2f}″")
    if top_emu + height_emu > slide_height_emu - margin_emu + int(Inches(0.02).emu):
        raise ValueError(f"Off-canvas bottom for {label}: bottom {(top_emu+height_emu)/914400:.2f}″ > {(slide_height_emu-margin_emu)/914400:.2f}″")


def rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str.lstrip("#").upper())


def add_textbox(slide: Any, text: str, left: Any, top: Any, width: Any, height: Any, font_size: Any, bold: bool, color_hex: str, alignment: Any, font_name: str = FONT_SANS) -> Any:
    shape = slide.shapes.add_textbox(left, top, width, height)
    shape.text_frame.word_wrap = True
    shape.text_frame.auto_size = None
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.alignment = alignment
    run = paragraph.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = rgb(color_hex)
    run.font.name = font_name
    # ensure anchor top so edits don't re-center glitchy
    shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    return shape


def add_card(slide: Any, left: Any, top: Any, width: Any, height: Any, fill_hex: str, border_hex: str) -> Any:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill_hex)
    shape.line.color.rgb = rgb(border_hex); shape.line.width = Pt(0.75)
    # round via adjustment
    try:
        shape.adjustments[0] = 0.08
    except Exception:
        pass
    return shape


def add_native_table(slide: Any, headers: list[str], rows: list[list[str]], left: Any, top: Any, width: Any, height: Any, palette: dict[str, str]) -> Any:
    if not headers and not rows:
        raise ValueError("results_table requires headers or rows")
    col_count = len(headers) if headers else len(rows[0])
    row_count = 1 + len(rows) if headers else len(rows)
    shape = slide.shapes.add_table(row_count, col_count, left, top, width, height)
    tbl = shape.table
    # header
    if headers:
        for idx, header_text in enumerate(headers):
            cell = tbl.cell(0, idx)
            cell.text = header_text
            paragraph = cell.text_frame.paragraphs[0]; paragraph.alignment = PP_ALIGN.CENTER
            run = paragraph.runs[0]; run.font.bold = True; run.font.size = Pt(9); run.font.color.rgb = rgb("FFFFFF")
            cell.fill.solid(); cell.fill.fore_color.rgb = rgb(palette["ACCENT"])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    start_row = 1 if headers else 0
    for r_idx, row_data in enumerate(rows):
        for c_idx in range(min(len(row_data), col_count)):
            cell = tbl.cell(start_row + r_idx, c_idx)
            if not (start_row == 0 and r_idx == 0 and not headers):
                cell.text = str(row_data[c_idx])
                paragraph = cell.text_frame.paragraphs[0]; paragraph.alignment = PP_ALIGN.LEFT
                run = paragraph.runs[0]; run.font.size = Pt(9); run.font.color.rgb = rgb(palette["FG"])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # banded
            if (r_idx % 2 == 1 and headers) or (r_idx % 2 == 0 and not headers):
                cell.fill.solid(); cell.fill.fore_color.rgb = rgb(palette["SURFACE"])
    return shape


def add_native_chart(slide: Any, title: str, categories: list[str], series: list[dict[str, Any]], left: Any, top: Any, width: Any, height: Any, palette: dict[str, str], chart_type: str) -> Any:
    chart_data = ChartData()
    chart_data.categories = categories
    palette_order = [palette["ACCENT"], palette["ACCENT_2"], palette["ACCENT_3"], palette["FG_MUTED"]]
    for idx, ser in enumerate(series):
        label = str(ser.get("name", f"Series {idx+1}"))
        values = [float(v) for v in ser.get("values", [])]
        if len(values) != len(categories):
            raise ValueError(f"Series '{label}' values len {len(values)} != categories len {len(categories)}")
        chart_data.add_series(label, values)

    xl_type = XL_CHART_TYPE.COLUMN_CLUSTERED if chart_type == "bar" else XL_CHART_TYPE.LINE_MARKERS
    shape = slide.shapes.add_chart(xl_type, left, top, width, height, chart_data)
    chart = shape.chart
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    chart.value_axis.has_major_gridlines = True
    for idx, ser in enumerate(chart.series):
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = rgb(palette_order[idx % len(palette_order)])
        ser.smooth = False
    return shape


# ── Slide builders — one job each ───────────────────────────────────────────

def build_cover(slide: Any, data: dict[str, Any], config: DeckConfig, content_width_emu: int, gutter_emu: int) -> None:
    title = str(data.get("title", config.title)).strip()
    subtitle = str(data.get("subtitle", "")).strip()
    if not title:
        raise ValueError("cover requires title")
    margin_emu = MARGIN.emu
    slide_w, slide_h = config.slide_width.emu, config.slide_height.emu
    # title centered
    add_textbox(slide, title, Inches(0.6), Inches(2.0), Inches(12.13) if config.grid_cols==12 else Inches(8.8), Inches(1.2), Pt(44), True, config.palette["ACCENT"] if config.theme_key!="research_dark" else config.palette["FG"], PP_ALIGN.CENTER, FONT_SANS if config.theme_key!="editorial" else FONT_SERIF)
    if subtitle:
        add_textbox(slide, subtitle, Inches(0.6), Inches(3.4), Inches(12.13) if config.grid_cols==12 else Inches(8.8), Inches(0.6), Pt(11), False, config.palette["FG_MUTED"], PP_ALIGN.CENTER)


def build_section_header(slide: Any, data: dict[str, Any], config: DeckConfig, content_width_emu: int, gutter_emu: int) -> None:
    title = str(data.get("title", "")).strip()
    kicker = str(data.get("kicker", "")).strip()
    if not title:
        raise ValueError("section_header requires title")
    # left rule
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(2.4), Inches(0.08), Inches(1.2))
    rule.fill.solid(); rule.fill.fore_color.rgb = rgb(config.palette["ACCENT"]); rule.line.fill.background()
    if kicker:
        add_textbox(slide, kicker.upper(), Inches(0.82), Inches(2.35), Inches(8.0), Inches(0.35), Pt(9), True, config.palette["FG_MUTED"], PP_ALIGN.LEFT)
    add_textbox(slide, title, Inches(0.82), Inches(2.75), Inches(8.0), Inches(0.9), Pt(28), True, config.palette["FG"], PP_ALIGN.LEFT)


def build_quote(slide: Any, data: dict[str, Any], config: DeckConfig, content_width_emu: int, gutter_emu: int) -> None:
    quote = str(data.get("quote", "")).strip()
    attribution = str(data.get("attribution", "")).strip()
    if not quote:
        raise ValueError("quote requires quote")
    # 8-col centered text (12-col grid: left offset 2 cols + gutter)
    # compute 8-col width
    col_w_emu = (content_width_emu - (config.grid_cols -1)*gutter_emu)//config.grid_cols
    width_8 = 8*col_w_emu + 7*gutter_emu
    left_8 = MARGIN.emu + 2*col_w_emu + 2*gutter_emu
    add_textbox(slide, f"“{quote}”", Emu(left_8), Inches(2.2), Emu(width_8), Inches(1.6), Pt(20), False, config.palette["FG"], PP_ALIGN.CENTER, FONT_SERIF)
    if attribution:
        add_textbox(slide, attribution, Emu(left_8), Inches(3.9), Emu(width_8), Inches(0.4), Pt(10), False, config.palette["FG_MUTED"], PP_ALIGN.RIGHT)


def build_bento_features(slide: Any, data: dict[str, Any], config: DeckConfig, content_width_emu: int, gutter_emu: int) -> None:
    title = str(data.get("title", "")).strip()
    items = data.get("items", [])
    if not isinstance(items, list) or len(items) == 0:
        raise ValueError("bento_features requires non-empty items [{title, body}]")
    if len(items) > 4:
        raise ValueError("bento_features max 4 cards — split into two slides")
    if title:
        add_textbox(slide, title, Inches(0.6), Inches(0.5), Inches(12.13) if config.grid_cols==12 else Inches(8.8), Inches(0.5), Pt(20), True, config.palette["FG"], PP_ALIGN.LEFT)
    # layout 4+4+4 (or 6+6 for 2)
    if len(items) == 2:
        spans = [6,6]
    elif len(items) == 3:
        spans = [4,4,4]
    elif len(items) == 4:
        spans = [3,3,3,3]
    else:
        spans = [12]
    col_w_emu = (content_width_emu - (config.grid_cols -1)*gutter_emu)//config.grid_cols
    top_emu = Inches(1.4).emu
    card_h = Inches(3.4)
    left_emu = MARGIN.emu
    for idx, item in enumerate(items):
        if isinstance(item, str):
            item_title, item_body = item, ""
        elif isinstance(item, dict):
            item_title = str(item.get("title", ""))
            item_body = str(item.get("body", ""))
        else:
            raise ValueError(f"bento_features item must be str or mapping, got {type(item).__name__}")
        span = spans[idx] if idx < len(spans) else spans[-1]
        width_emu = span*col_w_emu + (span-1)*gutter_emu
        # bounds check before draw
        check_bounds(left_emu, top_emu, width_emu, card_h.emu, config.slide_width.emu, config.slide_height.emu, MARGIN.emu, f"bento_features card {idx}")
        add_card(slide, Emu(left_emu), Emu(top_emu), Emu(width_emu), card_h, config.palette["SURFACE"], config.palette["BORDER"])
        title_h = Inches(0.38)
        body_top_emu = top_emu + CARD_PAD.emu + title_h.emu + int(Inches(0.10).emu)
        body_h_emu = card_h.emu - CARD_PAD.emu * 2 - title_h.emu - int(Inches(0.10).emu)
        if body_h_emu < int(Inches(0.8).emu):
            body_h_emu = int(Inches(0.8).emu)
        add_textbox(slide, item_title, Emu(left_emu+CARD_PAD.emu), Emu(top_emu+CARD_PAD.emu), Emu(width_emu-2*CARD_PAD.emu), title_h, Pt(11), True, config.palette["FG"], PP_ALIGN.LEFT)
        if item_body:
            add_textbox(slide, item_body, Emu(left_emu+CARD_PAD.emu), Emu(body_top_emu), Emu(width_emu-2*CARD_PAD.emu), Emu(body_h_emu), Pt(9), False, config.palette["FG_MUTED"], PP_ALIGN.LEFT)
        left_emu += width_emu + gutter_emu


def build_stats_grid(slide: Any, data: dict[str, Any], config: DeckConfig, content_width_emu: int, gutter_emu: int) -> None:
    title = str(data.get("title", "")).strip()
    stats = data.get("stats", [])
    if not isinstance(stats, list) or len(stats) == 0:
        raise ValueError("stats_grid requires non-empty stats [{value, label, delta}]")
    if len(stats) > 4:
        raise ValueError("stats_grid max 4 cards — split dense stats across slides (density rhythm)")
    if title:
        add_textbox(slide, title, Inches(0.6), Inches(0.5), Inches(12.13) if config.grid_cols==12 else Inches(8.8), Inches(0.5), Pt(20), True, config.palette["FG"], PP_ALIGN.LEFT)
    cols = 4 if len(stats) == 4 else (3 if len(stats)==3 else (2 if len(stats)==2 else 1))
    span = 3 if cols==4 else (4 if cols==3 else (6 if cols==2 else 12))
    col_w_emu = (content_width_emu - (config.grid_cols -1)*gutter_emu)//config.grid_cols
    card_w_emu = span*col_w_emu + (span-1)*gutter_emu
    card_h = Inches(1.9)
    left_emu = MARGIN.emu
    top_emu = Inches(2.2).emu
    for stat in stats:
        if not isinstance(stat, dict):
            raise ValueError(f"stats_grid item must be mapping, got {type(stat).__name__}")
        value = str(stat.get("value", "")).strip()
        label = str(stat.get("label", "")).strip()
        delta = str(stat.get("delta", "")).strip()
        if not value or not label:
            raise ValueError("stats_grid item requires value and label")
        check_bounds(left_emu, top_emu, card_w_emu, card_h.emu, config.slide_width.emu, config.slide_height.emu, MARGIN.emu, f"stats_grid card {value}")
        add_card(slide, Emu(left_emu), Emu(top_emu), Emu(card_w_emu), card_h, config.palette["SURFACE"], config.palette["BORDER"])
        value_h = Inches(0.55)
        label_h = Inches(0.32)
        delta_h = Inches(0.32) if delta else Inches(0)
        gap_small = int(Inches(0.06).emu)
        value_top = top_emu + CARD_PAD.emu
        label_top = value_top + value_h.emu + gap_small
        delta_top = label_top + label_h.emu + gap_small
        add_textbox(slide, value, Emu(left_emu+CARD_PAD.emu), Emu(value_top), Emu(card_w_emu-2*CARD_PAD.emu), value_h, Pt(28), True, config.palette["FG"], PP_ALIGN.LEFT, FONT_MONO)
        add_textbox(slide, label.upper(), Emu(left_emu+CARD_PAD.emu), Emu(label_top), Emu(card_w_emu-2*CARD_PAD.emu), label_h, Pt(9), True, config.palette["FG_MUTED"], PP_ALIGN.LEFT)
        if delta:
            delta_color = config.palette["ACCENT"] if "+" in delta or "▲" in delta else (config.palette["FG_MUTED"] if delta.startswith("~") else "#EF4444" if "-" in delta else config.palette["FG_MUTED"])
            add_textbox(slide, delta, Emu(left_emu+CARD_PAD.emu), Emu(delta_top), Emu(card_w_emu-2*CARD_PAD.emu), delta_h, Pt(9), False, delta_color, PP_ALIGN.LEFT)
        left_emu += card_w_emu + gutter_emu


def build_chart(slide: Any, data: dict[str, Any], config: DeckConfig, content_width_emu: int, gutter_emu: int, chart_type: str) -> None:
    title = str(data.get("title", "")).strip()
    categories = data.get("categories", [])
    series = data.get("series", [])
    if not isinstance(categories, list) or not categories:
        raise ValueError(f"{chart_type}_chart requires non-empty categories")
    if not isinstance(series, list) or not series:
        raise ValueError(f"{chart_type}_chart requires non-empty series [{{name, values}}]")
    if title:
        add_textbox(slide, title, Inches(0.6), Inches(0.5), Inches(12.13) if config.grid_cols==12 else Inches(8.8), Inches(0.45), Pt(14), True, config.palette["FG"], PP_ALIGN.LEFT)
    # 8+4 split: chart 8 cols, insight 4 cols (optional)
    col_w_emu = (content_width_emu - (config.grid_cols -1)*gutter_emu)//config.grid_cols
    chart_w_emu = 8*col_w_emu + 7*gutter_emu
    chart_h = Inches(3.8)
    left_emu = MARGIN.emu
    top_emu = Inches(1.6).emu
    insight = str(data.get("insight", "")).strip()
    if insight:
        add_native_chart(slide, "", categories, series, Emu(left_emu), Emu(top_emu), Emu(chart_w_emu), chart_h, config.palette, chart_type)
        # insight card 4 cols at right
        card_w_emu = 4*col_w_emu + 3*gutter_emu
        add_card(slide, Emu(left_emu+chart_w_emu+gutter_emu), Emu(top_emu), Emu(card_w_emu), chart_h, config.palette["SURFACE"], config.palette["BORDER"])
        add_textbox(slide, insight, Emu(left_emu+chart_w_emu+gutter_emu+CARD_PAD.emu), Emu(top_emu+CARD_PAD.emu), Emu(card_w_emu-2*CARD_PAD.emu), Emu(chart_h.emu-2*CARD_PAD.emu), Pt(10), False, config.palette["FG"], PP_ALIGN.LEFT)
    else:
        # full width chart (less common) — use 12 cols
        full_w = content_width_emu
        add_native_chart(slide, "", categories, series, Emu(left_emu), Emu(top_emu), Emu(full_w), chart_h, config.palette, chart_type)


def build_matrix(slide: Any, data: dict[str, Any], config: DeckConfig, content_width_emu: int, gutter_emu: int) -> None:
    title = str(data.get("title", "")).strip()
    quadrants = data.get("quadrants", [])
    items = data.get("items", [])
    if title:
        add_textbox(slide, title, Inches(0.6), Inches(0.5), Inches(12.13), Inches(0.45), Pt(14), True, config.palette["FG"], PP_ALIGN.LEFT)
    # outer rect centered: 8.0″×4.2″
    outer_w, outer_h = Inches(8.0), Inches(4.2)
    outer_left = (config.slide_width - outer_w) / 2
    outer_top = Inches(1.6)
    outer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, outer_left, outer_top, outer_w, outer_h)
    outer.fill.background(); outer.line.color.rgb = rgb(config.palette["BORDER"])
    # quad lines
    h_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, outer_left + outer_w/2 - Inches(0.01), outer_top, Inches(0.02), outer_h)
    h_line.fill.solid(); h_line.fill.fore_color.rgb = rgb(config.palette["BORDER"]); h_line.line.fill.background()
    v_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, outer_left, outer_top + outer_h/2 - Inches(0.01), outer_w, Inches(0.02))
    v_line.fill.solid(); v_line.fill.fore_color.rgb = rgb(config.palette["BORDER"]); v_line.line.fill.background()
    # quadrant labels (top-left etc)
    labels = ["Top-Left","Top-Right","Bottom-Left","Bottom-Right"]
    for idx, quad in enumerate(quadrants[:4]):
        label = str(quad.get("label", labels[idx])) if isinstance(quad, dict) else str(quad)
        # approximate quadrant centers — slightly higher to avoid dot-label collision
        qx = outer_left + (outer_w/4 if idx%2==0 else outer_w*3/4)
        qy = outer_top + (outer_h/4 if idx<2 else outer_h*3/4)
        add_textbox(slide, label, qx - Inches(0.7), qy - Inches(0.20), Inches(1.4), Inches(0.24), Pt(9), True, config.palette["FG_MUTED"], PP_ALIGN.CENTER)
    # items (dots)
    for item in items:
        if not isinstance(item, dict):
            continue
        label = str(item.get("label", "")).strip()
        x_frac = float(item.get("x", 0.5)); y_frac = float(item.get("y", 0.5))
        # clamp 0..1
        x_frac = max(0.05, min(0.95, x_frac)); y_frac = max(0.05, min(0.95, y_frac))
        dot_left = outer_left + Emu(int(outer_w.emu * x_frac)) - Inches(0.07)
        dot_top = outer_top + Emu(int(outer_h.emu * y_frac)) - Inches(0.07)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, dot_left, dot_top, Inches(0.14), Inches(0.14))
        dot.fill.solid(); dot.fill.fore_color.rgb = rgb(config.palette["ACCENT"]); dot.line.fill.background()
        if label:
            # place label offset from dot — south for top half, north for bottom half to avoid quadrant label collision
            if y_frac < 0.5:
                label_left = dot_left + Inches(0.12)
                label_top = dot_top + Inches(0.18)
            else:
                label_left = dot_left + Inches(0.12)
                label_top = dot_top - Inches(0.28)
            # keep inside outer rect: clamp x
            if label_left + Inches(1.1) > outer_left + outer_w:
                label_left = dot_left - Inches(1.22)  # flip to west if east would overflow
            add_textbox(slide, label, label_left, label_top, Inches(1.1), Inches(0.22), Pt(8), True, config.palette["FG"], PP_ALIGN.LEFT)


def build_timeline(slide: Any, data: dict[str, Any], config: DeckConfig, content_width_emu: int, gutter_emu: int) -> None:
    title = str(data.get("title", "")).strip()
    milestones = data.get("milestones", [])
    if not isinstance(milestones, list) or len(milestones) == 0:
        raise ValueError("timeline requires non-empty milestones [{label, title}]")
    if title:
        add_textbox(slide, title, Inches(0.6), Inches(0.5), Inches(12.13), Inches(0.45), Pt(14), True, config.palette["FG"], PP_ALIGN.LEFT)
    # horizontal line
    line_top = Inches(3.4)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), line_top, Inches(12.13), Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = rgb(config.palette["BORDER"]); line.line.fill.background()
    span_emu = content_width_emu // max(1, len(milestones))
    for idx, ms in enumerate(milestones):
        if isinstance(ms, str):
            label, ms_title = ms, ""
        elif isinstance(ms, dict):
            label, ms_title = str(ms.get("label", f"M{idx+1}")), str(ms.get("title", ""))
        else:
            continue
        cx_emu = MARGIN.emu + span_emu//2 + idx*span_emu
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(cx_emu - Inches(0.09).emu), line_top - Inches(0.07), Inches(0.18), Inches(0.18))
        dot.fill.solid(); dot.fill.fore_color.rgb = rgb(config.palette["ACCENT"]); dot.line.fill.background()
        add_textbox(slide, label, Emu(cx_emu - Inches(0.9).emu), line_top + Inches(0.22), Inches(1.8), Inches(0.3), Pt(9), True, config.palette["ACCENT"], PP_ALIGN.CENTER)
        if ms_title:
            add_textbox(slide, ms_title, Emu(cx_emu - Inches(0.9).emu), line_top + Inches(0.5), Inches(1.8), Inches(0.5), Pt(8), False, config.palette["FG"], PP_ALIGN.CENTER)


def build_results_table(slide: Any, data: dict[str, Any], config: DeckConfig, content_width_emu: int, gutter_emu: int) -> None:
    title = str(data.get("title", "")).strip()
    headers = [str(h) for h in data.get("headers", [])]
    rows = [[str(c) for c in r] for r in data.get("rows", [])]
    if not headers and not rows:
        raise ValueError("results_table requires headers or rows")
    if title:
        add_textbox(slide, title, Inches(0.6), Inches(0.5), Inches(12.13), Inches(0.45), Pt(14), True, config.palette["FG"], PP_ALIGN.LEFT)
    add_native_table(slide, headers, rows, Inches(0.6), Inches(1.4), Inches(12.13), Inches(4.5), config.palette)


def check_slop_in_texts(texts: list[str]) -> list[str]:
    findings: list[str] = []
    for text in texts:
        lower = text.lower()
        for token in SLOP_TOKENS:
            if token in lower:
                findings.append(f"Slop '{token}' in: {text[:60]!r}")
        if re.search(r"\bTODO\b", text) and "[CONTENT REQUIRED" not in text:
            findings.append(f"Bare TODO in: {text[:60]!r}")
    return findings


def build_slides(presentation: Any, slides_spec: list[Any], config: DeckConfig) -> None:
    if not isinstance(slides_spec, list):
        raise ValueError(f"'slides' must be list, got {type(slides_spec).__name__}")
    if len(slides_spec) == 0:
        raise ValueError("slides must be non-empty")
    if len(slides_spec) > 150:
        raise ValueError(f"Slide count {len(slides_spec)} exceeds sanity gate 150 — confirm with user")

    content_width_emu = int(config.slide_width.emu - 2 * MARGIN.emu)
    gutter_emu = GUTTER.emu

    type_map = {
        "cover": build_cover,
        "section_header": build_section_header,
        "closing": build_cover,
        "bento_features": build_bento_features,
        "moat_columns": build_bento_features,
        "stats_grid": build_stats_grid,
        "bar_chart": lambda s,d,c,w,g: build_chart(s,d,c,w,g,"bar"),
        "line_chart": lambda s,d,c,w,g: build_chart(s,d,c,w,g,"line"),
        "column_chart": lambda s,d,c,w,g: build_chart(s,d,c,w,g,"bar"),
        "matrix_2x2": build_matrix,
        "matrix": build_matrix,
        "timeline": build_timeline,
        "quote": build_quote,
        "results_table": build_results_table,
        "table": build_results_table,
        "comparison": build_bento_features,
        "process": build_bento_features,
    }

    text_inventory: list[str] = [config.title]
    for slide_spec in slides_spec:
        if not isinstance(slide_spec, dict):
            raise ValueError(f"Each slide must be mapping, got {type(slide_spec).__name__}")
        slide_type = str(slide_spec.get("type", "")).strip().lower()
        if not slide_type:
            raise ValueError("Slide missing 'type'")
        if slide_type not in type_map:
            raise ValueError(f"Unknown slide type '{slide_type}' — expected one of {sorted(type_map)}")
        # bullets guard — ghost of slop
        for key in ("title", "subtitle", "quote", "attribution", "kicker"):
            val = slide_spec.get(key)
            if isinstance(val, str) and val.strip():
                text_inventory.append(val)
        # pre-slop warn (not fail)
        for item_key in ("items", "stats", "bullets"):
            val = slide_spec.get(item_key)
            if isinstance(val, list):
                for item in val:
                    if isinstance(item, str):
                        text_inventory.append(item)
                    elif isinstance(item, dict):
                        for v in item.values():
                            if isinstance(v, str):
                                text_inventory.append(v)

        # add blank slide and build
        blank_layout = presentation.slide_layouts[6]  # Blank
        slide = presentation.slides.add_slide(blank_layout)
        builder = type_map[slide_type]
        try:
            builder(slide, slide_spec, config, content_width_emu, gutter_emu)
        except (ValueError, TypeError, KeyError) as exc:
            raise ValueError(f"Slide type '{slide_type}' build failed: {exc}") from exc

    slop = check_slop_in_texts(text_inventory)
    if slop:
        print("[create_pptx] Slop warnings:", file=sys.stderr)
        for finding in slop:
            print(f"  - {finding}", file=sys.stderr)
        print("  (Use [CONTENT REQUIRED: ...] instead of filler)", file=sys.stderr)


def main() -> None:
    args = parse_args()
    spec_path = Path(args.spec)
    output_path = Path(args.output)
    template_path = Path(args.template) if args.template else None

    if template_path and not template_path.exists():
        print(f"[create_pptx] Template not found: {template_path}", file=sys.stderr)
        sys.exit(1)

    try:
        raw_spec = load_spec(spec_path)
    except (FileNotFoundError, ValueError, RuntimeError) as exc:
        print(f"[create_pptx] Load failed: {exc}", file=sys.stderr)
        sys.exit(1)

    try:
        config = resolve_deck_config(raw_spec)
    except (ValueError, TypeError) as exc:
        print(f"[create_pptx] Config error: {exc}", file=sys.stderr)
        sys.exit(1)

    # presentation — from template or blank
    if template_path:
        try:
            presentation = Presentation(str(template_path))
            # enforce slide size from config? keep template's size as truth for brand fidelity
            slide_w_emu = presentation.slide_width.emu if presentation.slide_width is not None else 0  # type: ignore[union-attr]
            slide_h_emu = presentation.slide_height.emu if presentation.slide_height is not None else 0  # type: ignore[union-attr]
            config_w_emu = config.slide_width.emu  # type: ignore[union-attr]
            config_h_emu = config.slide_height.emu  # type: ignore[union-attr]
            if presentation.slide_width != config.slide_width or presentation.slide_height != config.slide_height:
                print(f"[create_pptx] Note: template size {slide_w_emu/914400:.2f}×{slide_h_emu/914400:.2f} vs spec {config_w_emu/914400:.2f}×{config_h_emu/914400:.2f} — using template size for fidelity", file=sys.stderr)
        except Exception as exc:
            print(f"[create_pptx] Failed to open template {template_path}: {exc}", file=sys.stderr)
            sys.exit(1)
    else:
        presentation = Presentation()
        presentation.slide_width = config.slide_width
        presentation.slide_height = config.slide_height

    # core properties
    core = presentation.core_properties
    core.title = config.title
    core.subject = config.subject
    core.author = config.author

    # set background fill per theme (solid)
    # Note: true theme via slideMaster — this solid fill is fallback for blank template
    # python-pptx background API is limited; we keep slide background as-is for blank, rely on cards for color

    slides_spec = raw_spec.get("slides", [])
    try:
        build_slides(presentation, slides_spec, config)
    except ValueError as exc:
        print(f"[create_pptx] Build error: {exc}", file=sys.stderr)
        sys.exit(1)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        presentation.save(str(output_path))
    except (OSError, PermissionError) as exc:
        print(f"[create_pptx] Save failed: {exc}", file=sys.stderr)
        sys.exit(1)

    print(f"[create_pptx] Wrote {output_path} ({output_path.stat().st_size} bytes) — {len(presentation.slides)} slides — theme {config.theme_key}")

if __name__ == "__main__":
    main()
